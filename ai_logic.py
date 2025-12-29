# ai_logic.py

import pdfplumber
from pptx import Presentation
import io
from fastapi import HTTPException

# Import the actual AI model function from your Model.py file
# Ensure Model.py is in the same directory
try:
    from Model import generate_explanation_and_example
except ImportError:
    print("Warning: Model.py not found or libraries missing. Ensure Model.py is present.")
    # Fallback function if Model.py fails to load (for testing purposes)
    def generate_explanation_and_example(text):
        return "Model not loaded.", "Please check server logs."

# ---------------------------------------------------------
# Helper function: Extract text from PDF
# Returns a list of dictionaries: [{"slide_number": 1, "text": "..."}]
# ---------------------------------------------------------
async def extract_pages_from_pdf(file_stream):
    pages_content = []
    try:
        # pdfplumber opens files directly
        with pdfplumber.open(file_stream) as pdf:
            for i, page in enumerate(pdf.pages):
                text = page.extract_text()
                # Only add pages that have text
                if text and text.strip():
                    pages_content.append({
                        "slide_number": i + 1,
                        "text": text.strip()
                    })
        return pages_content
    except Exception as e:
        print(f"Error extracting PDF: {e}")
        raise HTTPException(status_code=500, detail="Error occurred while parsing the PDF file.")

# ---------------------------------------------------------
# Helper function: Extract text from PPTX
# Returns a list of dictionaries: [{"slide_number": 1, "text": "..."}]
# ---------------------------------------------------------
async def extract_pages_from_pptx(file_stream):
    slides_content = []
    try:
        presentation = Presentation(file_stream)
        for i, slide in enumerate(presentation.slides):
            text = ""
            # Loop through shapes in the slide to find text frames
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
            
            # Only add slides that have text
            if text.strip():
                slides_content.append({
                    "slide_number": i + 1,
                    "text": text.strip()
                })
        return slides_content
    except Exception as e:
        print(f"Error extracting PPTX: {e}")
        raise HTTPException(status_code=500, detail="Error occurred while parsing the PPTX file.")

# ---------------------------------------------------------
# STEP 1: Process file for the 'Upload' screen
# Extracts text page by page and returns it to the frontend
# ---------------------------------------------------------
async def process_file_to_pages(file: io.BytesIO, file_type: str, filename: str):
    if file_type == "application/pdf":
        return await extract_pages_from_pdf(file)
    elif "presentation" in file_type or file_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        return await extract_pages_from_pptx(file)
    else:
        # Reject unsupported file types
        raise HTTPException(status_code=400, detail="Unsupported file type. Please upload a PDF or PPTX file.")

# ---------------------------------------------------------
# STEP 2: Process specific text for the 'Clarify' button
# Receives raw text from the frontend and sends it to the AI model
# ---------------------------------------------------------
async def process_text_directly(text: str):
    
    if not text.strip():
        return {
            "analysis": "No text provided.",
            "examples": ["Please select a slide containing text."]
        }

    try:
        # Call the function from Model.py
        explanation, example = generate_explanation_and_example(text)
    except Exception as e:
        print(f"AI Model Error: {e}")
        explanation = "Error generating explanation."
        example = str(e)

    # Construct the final response structure
    return {
        "analysis": explanation,
        "examples": [example] # The frontend expects a list
    }