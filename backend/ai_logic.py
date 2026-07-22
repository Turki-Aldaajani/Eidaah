# ai_logic.py

import pdfplumber
from pptx import Presentation
import io
from fastapi import HTTPException
import traceback

from concurrency import run_blocking

# Import the AI model function from Model.py
# Signature: generate_explanation_and_example(text: str) -> (explanation, example)
try:
    from Model import generate_explanation_and_example
    print("✅ Model.py imported successfully!")
except ImportError as e:
    print(f"❌ Warning: Model.py not found or libraries missing: {e}")
    traceback.print_exc()
    def generate_explanation_and_example(text):
        return "Model not loaded.", "Please check server logs."
except Exception as e:
    print(f"❌ Error loading Model.py: {e}")
    traceback.print_exc()
    def generate_explanation_and_example(text):
        return f"Model loading error: {str(e)}", "Check logs"


# ---------------------------------------------------------
# Constants
# ---------------------------------------------------------
MAX_FILE_SIZE_MB = 20
MAX_FILE_SIZE_BYTES = MAX_FILE_SIZE_MB * 1024 * 1024
ALLOWED_PDF_TYPE = "application/pdf"
ALLOWED_PPTX_TYPES = [
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "application/vnd.ms-powerpoint",
]
MAX_SLIDES = 100  # Safety limit
MAX_ANALYZE_CHARS = 6000  # Per-slide char cap for the Explain/analyze call


# ---------------------------------------------------------
# Validation
# ---------------------------------------------------------
def validate_file(file_stream: io.BytesIO, file_type: str, filename: str):
    """Validate file before processing."""
    # Check file size
    file_stream.seek(0, 2)  # Seek to end
    file_size = file_stream.tell()
    file_stream.seek(0)  # Reset to beginning

    if file_size == 0:
        raise HTTPException(status_code=400, detail="The uploaded file is empty.")

    if file_size > MAX_FILE_SIZE_BYTES:
        raise HTTPException(
            status_code=413,
            detail=f"File too large. Maximum size is {MAX_FILE_SIZE_MB}MB."
        )

    # Check file type
    is_pdf = file_type == ALLOWED_PDF_TYPE
    is_pptx = file_type in ALLOWED_PPTX_TYPES or "presentation" in (file_type or "")
    is_pptx_by_name = filename and filename.lower().endswith(".pptx")

    if not (is_pdf or is_pptx or is_pptx_by_name):
        raise HTTPException(
            status_code=400,
            detail="Unsupported file type. Please upload a PDF or PPTX file."
        )


# ---------------------------------------------------------
# Helper function: Extract text from PDF
# Returns a list of dictionaries: [{"slide_number": 1, "text": "..."}]
# ---------------------------------------------------------
def _extract_pages_from_pdf_sync(file_stream):
    pages_content = []
    try:
        with pdfplumber.open(file_stream) as pdf:
            for i, page in enumerate(pdf.pages):
                if i >= MAX_SLIDES:
                    break
                text = page.extract_text()
                if text and text.strip():
                    pages_content.append({
                        "slide_number": i + 1,
                        "text": text.strip()
                    })
        if not pages_content:
            raise HTTPException(
                status_code=400,
                detail="No readable text found in the PDF. The file may contain only images."
            )
        return pages_content
    except HTTPException:
        raise
    except Exception as e:
        print(f"Error extracting PDF: {e}")
        raise HTTPException(status_code=500, detail="Error occurred while parsing the PDF file.")


async def extract_pages_from_pdf(file_stream):
    """Parse the PDF on the shared pool — pdfplumber is fully blocking."""
    return await run_blocking(_extract_pages_from_pdf_sync, file_stream)


# ---------------------------------------------------------
# Helper function: Extract text from PPTX
# Returns a list of dictionaries: [{"slide_number": 1, "text": "..."}]
# ---------------------------------------------------------
def _extract_pages_from_pptx_sync(file_stream):
    slides_content = []
    try:
        presentation = Presentation(file_stream)
        for i, slide in enumerate(presentation.slides):
            if i >= MAX_SLIDES:
                break
            text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"

            if text.strip():
                slides_content.append({
                    "slide_number": i + 1,
                    "text": text.strip()
                })
        if not slides_content:
            raise HTTPException(
                status_code=400,
                detail="No readable text found in the presentation. Slides may contain only images."
            )
        return slides_content
    except HTTPException:
        raise
    except Exception as e:
        print(f"Error extracting PPTX: {e}")
        raise HTTPException(status_code=500, detail="Error occurred while parsing the PPTX file.")


async def extract_pages_from_pptx(file_stream):
    """Parse the PPTX on the shared pool — python-pptx is fully blocking."""
    return await run_blocking(_extract_pages_from_pptx_sync, file_stream)


# ---------------------------------------------------------
# STEP 1: Process file for the 'Upload' screen
# Extracts text page by page and returns it to the frontend
# ---------------------------------------------------------
async def process_file_to_pages(file: io.BytesIO, file_type: str, filename: str):
    # Validate first
    validate_file(file, file_type, filename)

    if file_type == ALLOWED_PDF_TYPE:
        return await extract_pages_from_pdf(file)
    else:
        # PPTX (already validated above)
        return await extract_pages_from_pptx(file)


# ---------------------------------------------------------
# STEP 2: Process specific text for the 'Clarify' button
# Receives raw text from the frontend and sends it to the AI model
# ---------------------------------------------------------
async def process_text_directly(text: str, language: str = None):
    print(f"\n{'='*60}")
    print(f"📝 Analyzing text (length: {len(text)} chars)")
    print(f"{'='*60}")

    if not text.strip():
        return {
            "analysis": "No text provided.",
            "examples": ["Please select a slide containing text."]
        }

    # Truncate very long texts to avoid excessive API usage.
    # Limit raised from 3000 to give more headroom; if we still have to cut,
    # tell the user via a `warning` in the response instead of failing silently.
    warning = None
    if len(text) > MAX_ANALYZE_CHARS:
        truncated_text = text[:MAX_ANALYZE_CHARS]
        print(f"⚠️  Text truncated from {len(text)} to {MAX_ANALYZE_CHARS} chars")
        warning = (
            f"تم تحليل أول {MAX_ANALYZE_CHARS} حرف فقط من هذه الشريحة لأنها طويلة جداً؛ قد لا يشمل الشرح كامل المحتوى."
            if language == "ar"
            else f"Only the first {MAX_ANALYZE_CHARS} characters of this slide were analyzed because it is very long; the explanation may not cover all of the content."
        )
    else:
        truncated_text = text

    try:
        print("🤖 Calling AI model...")
        explanation, example = await run_blocking(
            generate_explanation_and_example, truncated_text, language=language
        )
        print(f"✅ Analysis complete!")
        print(f"   Explanation: {explanation[:100]}...")
        print(f"   Example: {example[:100]}...")
    except Exception as e:
        print(f"❌ AI Model Error: {e}")
        traceback.print_exc()
        explanation = "An error occurred while analyzing this slide. Please try again."
        example = ""

    # Construct the final response structure (same shape as before)
    result = {
        "analysis": explanation,
        "examples": [example] if example else []
    }
    if warning:
        result["warning"] = warning

    print(f"{'='*60}\n")
    return result
