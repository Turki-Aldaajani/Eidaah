# slide_renderer.py
# Converts PDF pages and PPTX slides to JPEG images for visual preview.
#
# PDF:  pymupdf (fitz) — pure Python, no system dependencies
# PPTX: python-pptx → renders each slide to a PIL image → JPEG
#
# Both paths produce the same output: numbered JPEG files in a directory.

import os
import io
from pathlib import Path

# pymupdf (fitz) — pure Python PDF renderer
try:
    import fitz  # PyMuPDF
    FITZ_AVAILABLE = True
except ImportError:
    FITZ_AVAILABLE = False
    print("⚠️  pymupdf not installed — PDF slide preview disabled. Run: pip install pymupdf")

# python-pptx for PPTX rendering
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("⚠️  python-pptx not installed — PPTX slide preview disabled.")

# Pillow for image operations
try:
    from PIL import Image, ImageDraw, ImageFont
    PIL_AVAILABLE = True
except ImportError:
    PIL_AVAILABLE = False
    print("⚠️  Pillow not installed — PPTX slide preview disabled. Run: pip install Pillow")


# ---------------------
# Configuration
# ---------------------
RENDER_DPI = 150        # Balance between quality and file size
JPEG_QUALITY = 80       # JPEG compression quality (1–100)
SLIDE_WIDTH = 1280      # Text-rendered slide width (px)
SLIDE_HEIGHT = 720      # Text-rendered slide height (px)


# ---------------------
# PDF → JPEG via pymupdf
# ---------------------
def render_pdf_to_images(file_bytes: bytes, output_dir: str) -> list:
    """
    Render each PDF page as a JPEG image using pymupdf.
    Returns list of image filenames: ["1.jpg", "2.jpg", ...]
    """
    if not FITZ_AVAILABLE:
        return []

    os.makedirs(output_dir, exist_ok=True)

    try:
        doc = fitz.open(stream=file_bytes, filetype="pdf")
    except Exception as e:
        print(f"❌ PDF open error: {e}")
        return []

    filenames = []
    zoom = RENDER_DPI / 72  # pymupdf default is 72 DPI
    mat = fitz.Matrix(zoom, zoom)

    for page_num in range(len(doc)):
        try:
            page = doc.load_page(page_num)
            pix = page.get_pixmap(matrix=mat)
            filename = f"{page_num + 1}.jpg"
            filepath = os.path.join(output_dir, filename)
            pix.save(filepath)
            filenames.append(filename)
        except Exception as e:
            print(f"❌ Error rendering PDF page {page_num + 1}: {e}")

    doc.close()
    print(f"✅ Rendered {len(filenames)} PDF pages to images.")
    return filenames


# ---------------------
# PPTX → JPEG via python-pptx + Pillow
# ---------------------
def _get_shape_text(shape) -> str:
    """Safely extract text from a pptx shape."""
    try:
        if hasattr(shape, "text") and shape.text.strip():
            return shape.text.strip()
    except Exception:
        pass
    return ""


def _render_slide_as_image(slide, slide_index: int) -> "Image.Image | None":
    """
    Render a single pptx slide as a PIL Image using text layout.
    Returns None if Pillow is not available.
    """
    if not PIL_AVAILABLE:
        return None

    # Background
    img = Image.new("RGB", (SLIDE_WIDTH, SLIDE_HEIGHT), color=(15, 23, 42))  # dark blue-black
    draw = ImageDraw.Draw(img)

    # Slide number badge (top-left)
    badge_text = f"Slide {slide_index}"
    draw.rectangle([20, 20, 130, 50], fill=(59, 130, 246))  # blue badge
    draw.text((30, 27), badge_text, fill=(255, 255, 255))

    # Collect all text from slide shapes
    lines = []
    for shape in slide.shapes:
        text = _get_shape_text(shape)
        if text:
            lines.append(text)

    # Draw text content
    y = 80
    for line in lines:
        # Word-wrap manually: split into chunks of ~60 chars
        words = line.split()
        current = ""
        for word in words:
            if len(current) + len(word) + 1 <= 60:
                current = f"{current} {word}".strip()
            else:
                if current:
                    draw.text((40, y), current, fill=(224, 224, 224))
                    y += 28
                current = word
                if y > SLIDE_HEIGHT - 60:
                    draw.text((40, y), "...", fill=(128, 128, 128))
                    break
        if current and y <= SLIDE_HEIGHT - 60:
            draw.text((40, y), current, fill=(224, 224, 224))
            y += 28
        y += 10  # extra spacing between shapes

        if y > SLIDE_HEIGHT - 60:
            break

    # Bottom border accent
    draw.rectangle([0, SLIDE_HEIGHT - 4, SLIDE_WIDTH, SLIDE_HEIGHT], fill=(59, 130, 246))

    return img


def render_pptx_to_images(file_bytes: bytes, output_dir: str) -> list:
    """
    Render each PPTX slide as a JPEG image using python-pptx + Pillow.
    Returns list of image filenames: ["1.jpg", "2.jpg", ...]
    """
    if not PPTX_AVAILABLE or not PIL_AVAILABLE:
        return []

    os.makedirs(output_dir, exist_ok=True)

    try:
        prs = Presentation(io.BytesIO(file_bytes))
    except Exception as e:
        print(f"❌ PPTX open error: {e}")
        return []

    filenames = []
    for i, slide in enumerate(prs.slides):
        try:
            img = _render_slide_as_image(slide, i + 1)
            if img is None:
                continue
            filename = f"{i + 1}.jpg"
            filepath = os.path.join(output_dir, filename)
            img.save(filepath, "JPEG", quality=JPEG_QUALITY)
            filenames.append(filename)
        except Exception as e:
            print(f"❌ Error rendering PPTX slide {i + 1}: {e}")

    print(f"✅ Rendered {len(filenames)} PPTX slides to images.")
    return filenames


# ---------------------
# Main entry point
# ---------------------
def render_slides(file_bytes: bytes, file_type: str, filename: str, output_dir: str) -> list:
    """
    Detect file type and render slides to JPEG images.
    Returns list of image filenames (empty if rendering unavailable).
    """
    is_pdf = file_type == "application/pdf"
    is_pptx = (
        "presentation" in (file_type or "")
        or file_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        or (filename and filename.lower().endswith(".pptx"))
    )

    if is_pdf:
        return render_pdf_to_images(file_bytes, output_dir)
    elif is_pptx:
        return render_pptx_to_images(file_bytes, output_dir)
    else:
        print(f"⚠️  Unknown file type for rendering: {file_type}")
        return []
