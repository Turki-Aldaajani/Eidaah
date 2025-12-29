import os
import torch
import gradio as gr
from typing import List
from langchain_community.document_loaders import PyPDFLoader
from langchain_core.documents import Document
from pptx import Presentation
from difflib import SequenceMatcher
from transformers import pipeline, AutoTokenizer, AutoModelForCausalLM
import warnings

warnings.filterwarnings("ignore")

def gpu_usable():
    if not torch.cuda.is_available():
        return False
    try:
        x = torch.randn(1, device="cuda")
        y = x * 2
        torch.cuda.synchronize()
        return True
    except Exception as e:
        print(f"[GPU Disabled] CUDA present but unusable: {e}")
        return False


def get_llm():
    model_name = "Qwen/Qwen2-1.5B-Instruct"
    print(f"Loading optimized model: {model_name}")

    device = "cuda" if gpu_usable() else "cpu"
    print(f"Using device: {device}")

    tokenizer = AutoTokenizer.from_pretrained(model_name)

    model = AutoModelForCausalLM.from_pretrained(
        model_name,
        device_map={"": 0} if device == "cuda" else {"": "cpu"},
        torch_dtype=torch.float16 if device == "cuda" else torch.float32
    )

    pipe = pipeline(
        "text-generation",
        model=model,
        tokenizer=tokenizer,
        max_new_tokens=60,
        do_sample=False,
        return_full_text=False,
    )

    return pipe, tokenizer


llm_pipe, llm_tokenizer = get_llm()



# Universal Loader for PDF/PPTX
def load_document(file_path: str) -> List[Document]:
    ext = os.path.splitext(file_path)[1].lower()
    try:
        if ext == ".pdf":
            loader = PyPDFLoader(file_path)
            return loader.load()
        elif ext == ".pptx":
            prs = Presentation(file_path)
            slides = []
            for i, slide in enumerate(prs.slides):
                text = "\n".join(
                    shape.text for shape in slide.shapes if hasattr(shape, "text")
                )
                slides.append(Document(page_content=text.strip(), metadata={"page": i + 1}))
            return slides
        else:
            return []
    except Exception as e:
        print(f"Error loading document: {e}")
        return []

# (MAIN METHOD) Two-step generation
def generate_explanation_and_example(text: str):
    if not text.strip():
        return "No text found on this slide.", "No example available."
    try:
        # Step 1 — Explanation
        explain_prompt = f"Summarize this slide in two short sentences:\n\n{text}"
        messages = [{"role": "user", "content": explain_prompt}]
        formatted_prompt = llm_tokenizer.apply_chat_template(messages, tokenize=False, add_generation_prompt=True)
        explanation = llm_pipe(formatted_prompt, max_new_tokens=60)[0]["generated_text"].strip().lstrip(': \n')

        # Step 2 — Example
        example_prompt = f"""
Based ONLY on this explanation, give ONE short, clear real-life or classroom example
showing how it applies. Do NOT repeat the same sentences.

Explanation:
{explanation}
"""
        example_messages = [{"role": "user", "content": example_prompt}]
        formatted_example_prompt = llm_tokenizer.apply_chat_template(example_messages, tokenize=False, add_generation_prompt=True)
        example = llm_pipe(formatted_example_prompt, max_new_tokens=40)[0]["generated_text"].strip().lstrip(': \n')

        # Prevent copy-paste repetition
        similarity = SequenceMatcher(None, explanation.lower(), example.lower()).ratio()
        if similarity > 0.8:
            example = "Example: (could not generate a unique one). Try a different slide."

        return explanation, example
        
    except Exception as e:
        return f"Error: {e}", ""